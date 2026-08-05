import asyncio
from playwright.async_api import async_playwright
import os
from pptx import Presentation
from pptx.util import Inches

async def main():
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        # Use high resolution to get crisp images
        page = await browser.new_page(viewport={"width": 1920, "height": 1080})
        
        # Load local HTML file
        file_path = f"file:///{os.path.abspath('index.html').replace(chr(92), '/')}"
        await page.goto(file_path)
        
        # Wait a moment for fonts, particles, and images to load
        await page.wait_for_timeout(3000)
        
        # We need to hide controls for screenshots
        await page.evaluate("document.querySelector('.controls').style.display = 'none';")
        
        # Take screenshot of each slide
        slide_count = 8
        screenshots = []
        for i in range(1, slide_count + 1):
            # Make the current slide active, hide others
            await page.evaluate(f'''
                const slides = document.querySelectorAll('.slide');
                slides.forEach((slide, index) => {{
                    if (index === {i-1}) {{
                        slide.classList.add('active');
                        slide.style.opacity = '1';
                        slide.style.visibility = 'visible';
                        slide.style.transform = 'scale(1)';
                    }} else {{
                        slide.classList.remove('active');
                        slide.style.opacity = '0';
                        slide.style.visibility = 'hidden';
                        slide.style.transform = 'scale(0.95)';
                    }}
                }});
            ''')
            # Wait for any transition
            await page.wait_for_timeout(600)
            
            shot_path = f"assets/slide_{i}.png"
            await page.screenshot(path=shot_path)
            screenshots.append(shot_path)
            
        await browser.close()
        
        # Create PPTX
        prs = Presentation()
        prs.slide_width = Inches(13.333) # 16:9 for 1920x1080 proportion
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        for shot in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(shot, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save('StockSense_Presentation_Final.pptx')
        print("Successfully generated StockSense_Presentation_Final.pptx")

if __name__ == "__main__":
    asyncio.run(main())
