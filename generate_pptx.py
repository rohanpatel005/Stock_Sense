import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
# Change slide size to 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_COLOR = RGBColor(11, 11, 11)        # #0B0B0B
TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF
ACCENT_CYAN = RGBColor(0, 217, 255)    # #00D9FF
TEXT_GRAY = RGBColor(160, 160, 160)    # #A0A0A0
DANGER = RGBColor(255, 51, 102)        # #FF3366
SUCCESS = RGBColor(0, 255, 136)        # #00FF88

def set_slide_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text):
    title_shape = slide.shapes.title
    title_shape.text = text
    title_frame = title_shape.text_frame
    for p in title_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = 'Outfit'
            run.font.size = Pt(44)
            run.font.color.rgb = ACCENT_CYAN
            run.font.bold = True
    return title_shape

# Layouts
blank_layout = prs.slide_layouts[6]
title_layout = prs.slide_layouts[5]  # Title only

# --- Slide 1: Cover ---
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1)
txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "StockSense"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

p2 = tf.add_paragraph()
p2.text = "AI-Powered Stock Market Analysis & Paper Trading Platform\n\nFinal Year Project Presentation"
p2.font.size = Pt(24)
p2.font.color.rgb = TEXT_GRAY

if os.path.exists('assets/cover.png'):
    slide1.shapes.add_picture('assets/cover.png', Inches(7.5), Inches(1.5), width=Inches(5))

# --- Slide 2: Team ---
slide2 = prs.slides.add_slide(title_layout)
set_slide_bg(slide2)
add_title(slide2, "Team Members")

rows, cols = 4, 3
left = Inches(1)
top = Inches(2.5)
width = Inches(11)
height = Inches(2)
table = slide2.shapes.add_table(rows, cols, left, top, width, height).table

headers = ["Name", "Enrollment Number", "Roll Number"]
data = [
    ["Member 1 Name", "EN12345678", "RN001"],
    ["Member 2 Name", "EN12345679", "RN002"],
    ["Member 3 Name", "EN12345680", "RN003"]
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN
    cell.text_frame.paragraphs[0].font.size = Pt(20)

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx+1, col_idx)
        cell.text = val
        cell.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
        cell.text_frame.paragraphs[0].font.size = Pt(18)

# --- Slide 3: Problem vs Solution ---
slide3 = prs.slides.add_slide(title_layout)
set_slide_bg(slide3)
add_title(slide3, "Problem vs Solution")

problem_txBox = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(4))
tf_p = problem_txBox.text_frame
tf_p.word_wrap = True
p = tf_p.add_paragraph()
p.text = "The Problem"
p.font.color.rgb = DANGER
p.font.size = Pt(28)
p.font.bold = True
probs = ["Beginners lack stock market knowledge.", "Difficult to analyze stocks and technical indicators.", "Fear of losing money while investing.", "Research requires multiple websites and apps.", "No safe environment for beginners to practice."]
for prob in probs:
    p = tf_p.add_paragraph()
    p.text = "✖ " + prob
    p.font.color.rgb = TEXT_WHITE
    p.font.size = Pt(18)

solution_txBox = slide3.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(4))
tf_s = solution_txBox.text_frame
tf_s.word_wrap = True
p = tf_s.add_paragraph()
p.text = "The Solution (StockSense)"
p.font.color.rgb = SUCCESS
p.font.size = Pt(28)
p.font.bold = True
sols = ["AI Mentor explains investing in simple language.", "AI-powered stock analysis and technical indicators.", "Paper Trading with virtual balance.", "Portfolio tracking with AI Risk Detection.", "Live market data, news, and watchlist in one platform."]
for sol in sols:
    p = tf_s.add_paragraph()
    p.text = "✔ " + sol
    p.font.color.rgb = TEXT_WHITE
    p.font.size = Pt(18)

# --- Slide 4: Intro ---
slide4 = prs.slides.add_slide(title_layout)
set_slide_bg(slide4)
add_title(slide4, "Introduction to StockSense")

txBox = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(5.5), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "StockSense is an AI-powered stock market analysis and paper trading platform designed to help beginners and investors make informed investment decisions.\n\nIt combines real-time market data, AI-driven insights, technical analysis, portfolio tracking, and risk-free paper trading into a single platform, making investing simple, intelligent, and accessible."
p.font.size = Pt(22)
p.font.color.rgb = TEXT_WHITE

if os.path.exists('assets/intro.png'):
    slide4.shapes.add_picture('assets/intro.png', Inches(7), Inches(2), width=Inches(5.5))

# --- Slide 5: Tech Stack ---
slide5 = prs.slides.add_slide(title_layout)
set_slide_bg(slide5)
add_title(slide5, "Technology Stack")

techs = {
    "Frontend": "React.js, HTML5, CSS3, JavaScript, Tailwind CSS",
    "Backend": "Django, Django REST Framework, JWT Authentication",
    "Database & Auth": "PostgreSQL, Google OAuth, Email OTP",
    "APIs & Tools": "Gemini API, Yahoo Finance, Google News, GitHub"
}
top_margin = 1.8
for k, v in techs.items():
    txBox = slide5.shapes.add_textbox(Inches(1), Inches(top_margin), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = k
    p.font.color.rgb = ACCENT_CYAN
    p.font.size = Pt(24)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = v
    p2.font.color.rgb = TEXT_WHITE
    p2.font.size = Pt(18)
    top_margin += 1.3

# --- Slide 6: Workflow ---
slide6 = prs.slides.add_slide(title_layout)
set_slide_bg(slide6)
add_title(slide6, "StockSense Workflow")

txBox = slide6.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Login / Signup ➔ Dashboard ➔ Search Stock ➔ AI Analysis ➔ Paper Trading ➔ Portfolio Tracking ➔ Smart Investment Decision"
p.font.color.rgb = TEXT_WHITE
p.font.size = Pt(24)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# --- Slide 7: ER Diagram ---
slide7 = prs.slides.add_slide(title_layout)
set_slide_bg(slide7)
add_title(slide7, "Database Design (ER Diagram)")
txBox = slide7.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "[ Users ]\nuser_id (PK), full_name, email, password, wallet_balance, created_at\n\n[ Portfolio ]\nportfolio_id (PK), user_id (FK), stock_symbol, quantity, average_price\n\n[ Transactions ]\ntransaction_id (PK), user_id (FK), stock_symbol, type, quantity, price, date\n\n[ Watchlist ]\nwatchlist_id (PK), user_id (FK), stock_symbol"
p.font.color.rgb = TEXT_WHITE
p.font.size = Pt(18)

# --- Slide 8: Future Scope ---
slide8 = prs.slides.add_slide(title_layout)
set_slide_bg(slide8)
add_title(slide8, "Advantages & Future Scope")

adv_txBox = slide8.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(4))
tf_a = adv_txBox.text_frame
tf_a.word_wrap = True
p = tf_a.add_paragraph()
p.text = "Advantages"
p.font.color.rgb = ACCENT_CYAN
p.font.size = Pt(28)
p.font.bold = True
advs = ["AI-powered investment guidance", "Beginner-friendly interface", "Risk-free paper trading", "Real-time market insights", "Portfolio performance tracking", "AI Risk Detection", "Personalized watchlists"]
for a in advs:
    p = tf_a.add_paragraph()
    p.text = "• " + a
    p.font.color.rgb = TEXT_WHITE
    p.font.size = Pt(16)

fut_txBox = slide8.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(4))
tf_f = fut_txBox.text_frame
tf_f.word_wrap = True
p = tf_f.add_paragraph()
p.text = "Future Scope"
p.font.color.rgb = ACCENT_CYAN
p.font.size = Pt(28)
p.font.bold = True
futs = ["Live trading integration", "AI price prediction", "Voice-enabled AI assistant", "Multi-language support", "Advanced portfolio analytics", "Community discussion forum", "Mobile application"]
for f in futs:
    p = tf_f.add_paragraph()
    p.text = "• " + f
    p.font.color.rgb = TEXT_WHITE
    p.font.size = Pt(16)

if os.path.exists('assets/future.png'):
    slide8.shapes.add_picture('assets/future.png', Inches(10), Inches(2), width=Inches(3))

prs.save('StockSense_Presentation.pptx')
print("Successfully generated StockSense_Presentation.pptx")
